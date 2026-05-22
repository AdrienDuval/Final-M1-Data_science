#!/usr/bin/env python3
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

BG     = RGBColor(0xF0,0xF4,0xF8); WHITE  = RGBColor(0xFF,0xFF,0xFF)
ACCENT = RGBColor(0x25,0x63,0xEB); BLUE_D = RGBColor(0x1E,0x40,0xAF)
BLUE_L = RGBColor(0xDB,0xEA,0xFE); TEXT_P = RGBColor(0x0F,0x17,0x2A)
TEXT_S = RGBColor(0x33,0x41,0x55); TEXT_M = RGBColor(0x64,0x74,0x8B)
TV     = RGBColor(0x3B,0x82,0xF6); RADIO  = RGBColor(0x8B,0x5C,0xF6)
SOCIAL = RGBColor(0x06,0xB6,0xD4); GREEN  = RGBColor(0x10,0xB9,0x81)
AMBER  = RGBColor(0xF5,0x9E,0x0B); RED_C  = RGBColor(0xEF,0x44,0x44)
BORDER = RGBColor(0xE2,0xE8,0xF0); INDIGO = RGBColor(0x63,0x66,0xF1)
COVER_R= RGBColor(0x1E,0x40,0xAF); PANEL_DOT=RGBColor(0x17,0x32,0x89)

FONT = "Segoe UI"
W, H = Inches(13.33), Inches(7.5)

# ── primitives ───────────────────────────────────────────────────────────────

def bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def _s(slide, t, x, y, w, h, fill=None, line=None, lw=Pt(0.5)):
    s = slide.shapes.add_shape(t, x, y, w, h)
    if fill is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = lw
    return s

def rect(sl,x,y,w,h,fill=None,line=None,lw=Pt(0.5)):  return _s(sl,1,x,y,w,h,fill,line,lw)
def rrect(sl,x,y,w,h,fill=None,line=None,lw=Pt(0.5)): return _s(sl,5,x,y,w,h,fill,line,lw)
def oval(sl,x,y,w,h,fill=None):                        return _s(sl,9,x,y,w,h,fill)

def txt(sl, text, x, y, w, h, size=Pt(12), bold=False, color=TEXT_P,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = sl.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = size
    r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color
    return box

def stxt(shape, text, size=Pt(12), bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    tf = shape.text_frame; tf.word_wrap = True
    try:
        bp = tf._txBody.find(qn('a:bodyPr'))
        if bp is not None: bp.set('anchor', 'ctr')
    except Exception: pass
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run() if not p.runs else p.runs[0]
    r.text = text; r.font.name = FONT; r.font.size = size
    r.font.bold = bold; r.font.color.rgb = color

def circ(sl, x, y, size, label, color):
    s = oval(sl, x, y, size, size, fill=color); stxt(s, label, size=Pt(13), bold=True)
    return s

def new_slide(prs): return prs.slides.add_slide(prs.slide_layouts[6])

def header(sl, title, sub=None):
    rect(sl, 0, 0, W, Inches(0.09), fill=ACCENT)
    txt(sl, title, Inches(0.55), Inches(0.18), Inches(12.0), Inches(0.6),
        size=Pt(25), bold=True, color=TEXT_P)
    if sub:
        txt(sl, sub, Inches(0.55), Inches(0.76), Inches(11.0), Inches(0.4),
            size=Pt(11), color=TEXT_M)
    rect(sl, Inches(0.55), Inches(1.14), Inches(12.23), Inches(0.018), fill=BORDER)

# ── SLIDE 1 — Cover ─────────────────────────────────────────────────────────

def s1(prs):
    sl = new_slide(prs); bg(sl, BG)
    # Blue right panel
    rect(sl, Inches(8.0), 0, Inches(5.33), H, fill=COVER_R)
    # Decorative circles on panel
    dots = [(8.4,0.4,1.8),(9.8,1.1,0.7),(10.9,0.3,1.1),(11.7,1.8,0.5),
            (8.2,2.8,0.6),(9.0,3.5,1.4),(11.2,3.2,0.9),(10.3,4.6,1.6),
            (8.6,5.5,0.8),(11.5,5.0,1.2),(9.5,6.2,0.6),(12.0,6.5,0.7)]
    for (dx,dy,ds) in dots:
        oval(sl, Inches(dx), Inches(dy), Inches(ds), Inches(ds), fill=PANEL_DOT)
    # Highlight circle
    oval(sl, Inches(9.3), Inches(2.0), Inches(2.8), Inches(2.8), fill=RGBColor(0x1D,0x4E,0xD8))
    oval(sl, Inches(10.1), Inches(2.8), Inches(1.2), Inches(1.2), fill=RGBColor(0x93,0xC5,0xFD))
    # Names on panel
    txt(sl,"EFREI Paris · M1 Data Science & IA · 2025–2026",
        Inches(8.2),Inches(5.9),Inches(4.9),Inches(0.35),
        size=Pt(9),color=RGBColor(0xBF,0xDB,0xFE),align=PP_ALIGN.CENTER)
    txt(sl,"CHUEMBOU ADRIEN  ·  QUANG DAT",
        Inches(8.2),Inches(6.3),Inches(4.9),Inches(0.38),
        size=Pt(11),bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    # Left accent bar
    rect(sl, 0, 0, Inches(0.07), H, fill=ACCENT)
    # Title
    txt(sl,"ROI Intelligence",
        Inches(0.7),Inches(1.1),Inches(7.0),Inches(1.05),
        size=Pt(46),bold=True,color=TEXT_P)
    txt(sl,"Optimisation du retour sur investissement marketing\npar le Machine Learning",
        Inches(0.7),Inches(2.25),Inches(7.0),Inches(0.9),
        size=Pt(16),color=ACCENT)
    txt(sl,"Système de prédiction multi-modèles et tableau de bord décisionnel en temps réel",
        Inches(0.7),Inches(3.25),Inches(7.0),Inches(0.55),
        size=Pt(11.5),color=TEXT_S)
    rect(sl, Inches(0.7), Inches(3.95), Inches(4.2), Inches(0.025), fill=BORDER)
    # KPI chips
    chips = [("R² 0.9965","Score modèle",TV),("×4","Modèles ML",RADIO),
             ("×7","Pages dashboard",SOCIAL),("<50ms","Latence API",GREEN)]
    for i,(val,lbl,col) in enumerate(chips):
        cx = Inches(0.7)+i*Inches(1.75); cy = Inches(4.25)
        rrect(sl,cx,cy,Inches(1.62),Inches(0.72),fill=WHITE,line=BORDER)
        rect(sl,cx,cy,Inches(0.065),Inches(0.72),fill=col)
        txt(sl,val,cx+Inches(0.13),cy+Inches(0.05),Inches(1.35),Inches(0.32),
            size=Pt(14),bold=True,color=col)
        txt(sl,lbl,cx+Inches(0.13),cy+Inches(0.40),Inches(1.35),Inches(0.26),
            size=Pt(8.5),color=TEXT_M)
    return sl

# ── SLIDE 2 — Problématique ──────────────────────────────────────────────────

def s2(prs):
    sl = new_slide(prs); bg(sl, BG)
    header(sl,"Problématique",
           "Les budgets publicitaires sont encore trop souvent alloués à l'intuition — pas à la donnée.")
    # Quote banner
    rrect(sl,Inches(0.55),Inches(1.26),Inches(12.23),Inches(0.46),
          fill=BLUE_L,line=ACCENT,lw=Pt(1.0))
    txt(sl,"L'enjeu : transformer une intuition budgétaire en décision data-driven, mesurable et reproductible.",
        Inches(0.75),Inches(1.32),Inches(12.0),Inches(0.36),
        size=Pt(11),color=BLUE_D,italic=True)
    # Two columns
    cy,ch,cw = Inches(1.86),Inches(5.15),Inches(5.95)
    # Left card
    rrect(sl,Inches(0.55),cy,cw,ch,fill=WHITE,line=BORDER)
    rect(sl,Inches(0.55),cy,Inches(0.07),ch,fill=RED_C)
    txt(sl,"❌  Approche traditionnelle",
        Inches(0.75),cy+Inches(0.14),Inches(5.6),Inches(0.38),
        size=Pt(13),bold=True,color=RED_C)
    trad = ["Allocation empirique basée sur l'expérience passée",
            "Aucune corrélation mesurée entre budget et revenu",
            "Décisions prises a posteriori, après les campagnes",
            "Même budget, résultats imprévisibles selon les canaux",
            "Impossible de simuler l'impact avant d'engager"]
    for i,item in enumerate(trad):
        iy = cy+Inches(0.65)+i*Inches(0.78)
        oval(sl,Inches(0.78),iy+Inches(0.1),Inches(0.09),Inches(0.09),fill=RED_C)
        txt(sl,item,Inches(1.0),iy,Inches(5.3),Inches(0.65),size=Pt(10.5),color=TEXT_S,wrap=True)
    # Right card
    rrect(sl,Inches(6.83),cy,cw,ch,fill=WHITE,line=BORDER)
    rect(sl,Inches(6.83),cy,Inches(0.07),ch,fill=GREEN)
    txt(sl,"✅  Notre approche ML",
        Inches(7.03),cy+Inches(0.14),Inches(5.6),Inches(0.38),
        size=Pt(13),bold=True,color=GREEN)
    ml = ["Modèles entraînés sur 4 572 campagnes documentées",
          "Prédiction du CA en temps réel selon le budget alloué",
          "4 canaux : TV · Radio · Social Media · Influenceur",
          "R² > 0.994 — corrélation quasi-parfaite avec les ventes",
          "Simulation live : le revenu s'affiche instantanément"]
    for i,item in enumerate(ml):
        iy = cy+Inches(0.65)+i*Inches(0.78)
        oval(sl,Inches(7.06),iy+Inches(0.1),Inches(0.09),Inches(0.09),fill=GREEN)
        txt(sl,item,Inches(7.28),iy,Inches(5.3),Inches(0.65),size=Pt(10.5),color=TEXT_S,wrap=True)
    # VS badge
    vs = oval(sl,Inches(6.34),cy+Inches(2.2),Inches(0.65),Inches(0.65),fill=ACCENT)
    stxt(vs,"VS",size=Pt(11),bold=True)
    return sl

# ── SLIDE 3 — Objectifs ──────────────────────────────────────────────────────

def s3(prs):
    sl = new_slide(prs); bg(sl, BG)
    header(sl,"Objectifs du projet",
           "Quatre ambitions concrètes, de la donnée brute à la décision business éclairée.")
    objs = [
        ("①","Prédire",TV,
         "Construire 4 modèles ML capables de prédire le chiffre d'affaires\nà partir du mix budgétaire TV · Radio · Social Media · Influenceur."),
        ("②","Comparer & Sélectionner",RADIO,
         "Évaluer chaque modèle via validation croisée 5 plis (R², MAE, RMSE)\net sélectionner automatiquement le plus performant."),
        ("③","Exposer via API REST",SOCIAL,
         "Déployer les modèles entraînés via une API FastAPI scalable,\ninterrogeable en temps réel par n'importe quel client."),
        ("④","Visualiser & Décider",GREEN,
         "Proposer un dashboard interactif en 7 pages permettant\nde simuler, comparer et optimiser les allocations budgétaires."),
    ]
    cw,ch = Inches(5.9),Inches(2.48)
    pos = [(Inches(0.55),Inches(1.3)),(Inches(6.88),Inches(1.3)),
           (Inches(0.55),Inches(3.95)),(Inches(6.88),Inches(3.95))]
    for (num,title,col,desc),(cx,cy) in zip(objs,pos):
        rrect(sl,cx,cy,cw,ch,fill=WHITE,line=BORDER)
        rect(sl,cx,cy,Inches(0.07),ch,fill=col)
        circ(sl,cx+Inches(0.2),cy+Inches(0.17),Inches(0.46),num,col)
        txt(sl,title,cx+Inches(0.82),cy+Inches(0.19),Inches(4.9),Inches(0.4),
            size=Pt(14),bold=True,color=TEXT_P)
        rect(sl,cx+Inches(0.18),cy+Inches(0.72),cw-Inches(0.26),Inches(0.015),fill=BORDER)
        txt(sl,desc,cx+Inches(0.18),cy+Inches(0.84),cw-Inches(0.28),Inches(1.48),
            size=Pt(10.5),color=TEXT_S,wrap=True)
    return sl

# ── SLIDE 4 — Architecture ───────────────────────────────────────────────────

def s4(prs):
    sl = new_slide(prs); bg(sl, BG)
    header(sl,"Architecture de la solution",
           "Cinq couches modulaires — de la donnée brute à la décision business.")
    layers = [
        ("DONNÉES",        TV,    "data/Dummy Data HSS.csv",
         "4 572 lignes · TV, Radio, Social Media, Influencer, Sales"),
        ("PREPROCESSING",  SOCIAL,"src/preprocessing.py",
         "StandardScaler + OneHotEncoder · Pipeline sklearn · Imputation médiane"),
        ("ENTRAÎNEMENT ML",RADIO, "src/train.py",
         "RandomizedSearchCV · 15 iter · 5-fold CV · 4 algorithmes → .pkl (joblib)"),
        ("API REST",       ACCENT,"api/main.py — FastAPI",
         "/predict · /predict/all · /metrics · /feature-importance · /stats · /health"),
        ("DASHBOARD",      GREEN, "dashboard-next/ — Next.js 14",
         "TypeScript · Tailwind CSS · Recharts · Framer Motion · 7 pages · Clair/Sombre"),
    ]
    lh,lw,gap = Inches(0.9),Inches(12.23),Inches(0.06)
    sx,sy = Inches(0.55),Inches(1.3)
    for i,(name,col,path,desc) in enumerate(layers):
        y = sy+i*(lh+gap)
        rrect(sl,sx,y,lw,lh,fill=WHITE,line=BORDER)
        blk = rect(sl,sx,y,Inches(1.5),lh,fill=col)
        stxt(blk,name,size=Pt(9),bold=True,color=WHITE)
        txt(sl,path,sx+Inches(1.66),y+Inches(0.1),Inches(10.4),Inches(0.38),
            size=Pt(11.5),bold=True,color=TEXT_P)
        txt(sl,desc,sx+Inches(1.66),y+Inches(0.5),Inches(10.4),Inches(0.34),
            size=Pt(10),color=TEXT_M)
        if i<4:
            txt(sl,"▼",sx+Inches(5.6),y+lh,Inches(0.4),Inches(gap),
                size=Pt(7),color=TEXT_M,align=PP_ALIGN.CENTER)
    # docker command bar
    dy = sy+5*(lh+gap)+Inches(0.04)
    rrect(sl,sx,dy,lw,Inches(0.4),fill=TEXT_P)
    txt(sl,"⚡  docker-compose up  —  API :8000  ·  Dashboard :3000",
        sx+Inches(0.2),dy+Inches(0.06),Inches(12.0),Inches(0.3),
        size=Pt(11),bold=True,color=WHITE)
    return sl

# ── SLIDE 5 — Pipeline ML ────────────────────────────────────────────────────

def s5(prs):
    sl = new_slide(prs); bg(sl, BG)
    header(sl,"Pipeline ML : de la donnée brute à la prédiction",
           "Un pipeline rigoureusement validé — de la saisie brute à l'inférence en moins de 50 ms.")
    steps = [
        ("①","load_data()","CSV · 4 572 lignes\nValidation schéma",TV),
        ("②","build_preprocessor()","StandardScaler\n+ OneHotEncoder",RADIO),
        ("③","tune()","RandomizedSearchCV\n15 iter · 5-fold",SOCIAL),
        ("④","train_all()","4 modèles .pkl\nsplits.pkl",GREEN),
        ("⑤","FastAPI","Lifespan loading\n< 50ms / req.",ACCENT),
        ("⑥","Dashboard","Sliders → API\nTemps réel",INDIGO),
    ]
    sw,sh = Inches(1.82),Inches(4.6)
    sy,sx0 = Inches(1.75),Inches(0.55)
    unit = Inches(2.13)  # sw + arrow space
    for i,(num,func,desc,col) in enumerate(steps):
        sx = sx0+i*unit
        rrect(sl,sx,sy,sw,sh,fill=WHITE,line=BORDER)
        rect(sl,sx,sy,sw,Inches(0.065),fill=col)
        circ(sl,sx+(sw-Inches(0.46))/2,sy+Inches(0.12),Inches(0.46),num,col)
        txt(sl,func,sx+Inches(0.08),sy+Inches(0.72),sw-Inches(0.16),Inches(0.42),
            size=Pt(9.5),bold=True,color=TEXT_P,align=PP_ALIGN.CENTER)
        txt(sl,desc,sx+Inches(0.08),sy+Inches(1.22),sw-Inches(0.16),Inches(3.1),
            size=Pt(9.5),color=TEXT_M,align=PP_ALIGN.CENTER,wrap=True)
        if i<5:
            txt(sl,"›",sx+sw+Inches(0.12),sy+Inches(1.8),Inches(0.2),Inches(0.45),
                size=Pt(18),color=TEXT_M,align=PP_ALIGN.CENTER)
    # note
    ny = sy+sh+Inches(0.2)
    rrect(sl,sx0,ny,Inches(12.23),Inches(0.5),fill=BLUE_L,line=ACCENT,lw=Pt(0.75))
    txt(sl,"🔒  models/splits.pkl garantit la reproductibilité : métriques toujours calculées sur le même jeu de test.",
        sx0+Inches(0.2),ny+Inches(0.1),Inches(12.0),Inches(0.35),
        size=Pt(10.5),color=BLUE_D)
    return sl

# ── SLIDE 6 — Modèles ML ─────────────────────────────────────────────────────

def s6(prs):
    sl = new_slide(prs); bg(sl, BG)
    header(sl,"Les 4 modèles ML — Performances comparées",
           "Tous les modèles dépassent R² = 0.99 — le Random Forest s'impose comme référence.")
    models = [
        ("🏆 Random Forest",   0.9965,0.0028,TV,    True),
        ("Gradient Boosting",  0.9964,0.0028,RADIO,  False),
        ("Régression Linéaire",0.9956,0.0032,SOCIAL, False),
        ("MLP Neural Network", 0.9941,0.0034,INDIGO, False),
    ]
    # Table (left)
    tx,ty = Inches(0.55),Inches(1.32)
    cols = [Inches(3.1),Inches(1.5),Inches(1.3),Inches(1.0)]
    rh = Inches(0.52)
    rect(sl,tx,ty,sum(cols),rh,fill=ACCENT)
    for hi,(hl,cw) in enumerate(zip(["Modèle","R² moyen","Éc.-type","Statut"],cols)):
        cx = tx+sum(cols[:hi])
        txt(sl,hl,cx+Inches(0.1),ty+Inches(0.12),cw,rh,size=Pt(10.5),bold=True,color=WHITE)
    for ri,(name,r2,std,col,best) in enumerate(models):
        ry = ty+(ri+1)*rh
        rf = BLUE_L if best else (WHITE if ri%2==0 else RGBColor(0xF8,0xFA,0xFF))
        rect(sl,tx,ry,sum(cols),rh,fill=rf,line=BORDER,lw=Pt(0.3))
        vals = [name,f"{r2:.4f}",f"±{std:.4f}","BEST" if best else "—"]
        fc   = [ACCENT if best else TEXT_P, TV if best else TEXT_S,
                TEXT_M, ACCENT if best else TEXT_M]
        bolds= [best, False, False, best]
        for vi,(v,cw) in enumerate(zip(vals,cols)):
            cx = tx+sum(cols[:vi])
            txt(sl,v,cx+Inches(0.1),ry+Inches(0.1),cw,rh,
                size=Pt(10.5),bold=bolds[vi],color=fc[vi])
    # Bar chart (right)
    bx,by = Inches(7.1),Inches(1.32)
    bw,bh = Inches(5.7),Inches(4.6)
    rrect(sl,bx,by,bw,bh,fill=WHITE,line=BORDER)
    txt(sl,"R² par modèle  —  échelle 0.990 → 1.000",
        bx+Inches(0.2),by+Inches(0.14),Inches(5.3),Inches(0.3),
        size=Pt(9.5),color=TEXT_M,bold=True)
    bar_x = bx+Inches(1.65); max_bw = Inches(3.65)
    bnames = ["Random Forest","Grad. Boosting","Régression Lin.","MLP"]
    bcolors= [TV,RADIO,SOCIAL,INDIGO]
    bvals  = [0.9965,0.9964,0.9956,0.9941]
    for i,(bn,bc,bv) in enumerate(zip(bnames,bcolors,bvals)):
        iy = by+Inches(0.6)+i*Inches(0.92)
        bar_fill_w = ((bv-0.990)/0.010)*max_bw
        txt(sl,bn,bx+Inches(0.15),iy+Inches(0.1),Inches(1.45),Inches(0.35),
            size=Pt(9),color=TEXT_S)
        rect(sl,bar_x,iy+Inches(0.08),max_bw,Inches(0.35),
             fill=RGBColor(0xF1,0xF5,0xF9))
        rect(sl,bar_x,iy+Inches(0.08),bar_fill_w,Inches(0.35),fill=bc)
        txt(sl,f"{bv:.4f}",bar_x+bar_fill_w+Inches(0.08),iy+Inches(0.08),
            Inches(0.7),Inches(0.35),size=Pt(9.5),bold=True,color=bc)
    # Best model banner
    rrect(sl,Inches(0.55),Inches(6.22),Inches(12.23),Inches(0.75),fill=BLUE_L,line=ACCENT)
    txt(sl,"🏆  Meilleur modèle : Random Forest  —  R² = 0.9965  ·  Explique 99.65 % de la variance des ventes",
        Inches(0.75),Inches(6.38),Inches(12.0),Inches(0.44),
        size=Pt(12.5),bold=True,color=BLUE_D)
    return sl

# ── SLIDE 7 — Dashboard pages ────────────────────────────────────────────────

def s7(prs):
    sl = new_slide(prs); bg(sl, BG)
    header(sl,"Dashboard — 7 pages d'intelligence marketing",
           "De la prédiction brute à l'optimisation stratégique — chaque page répond à une question business précise.")
    pages = [
        ("📊","Accueil",         TV,    "Vue globale · KPIs · corrélation canaux · benchmark revenus"),
        ("🎯","Prévision",       RADIO, "Sliders · prédiction temps réel · ROI · consensus 4 modèles"),
        ("⚙️","Optimiseur",      SOCIAL,"Présets · sensibilité TV · comparaison des modèles"),
        ("🗺️","Planificateur",   GREEN, "Objectif → budget optimal · inversion modèle · Monte Carlo RF"),
        ("🤖","Comparaison ML",  INDIGO,"Métriques · diagramme radar · matrice de confusion"),
        ("💡","Insights",        AMBER, "Importance features · corrélation · ROI par influenceur"),
        ("📈","Feature Importan.",ACCENT,"Classement SHAP · permutation importance · stats canal-revenu"),
    ]
    # Top row: 4 cards, bottom row: 3 cards centered
    top_cw,bot_cw = Inches(2.98),Inches(3.78)
    card_h = Inches(2.7)
    top_y,bot_y = Inches(1.32),Inches(4.18)
    top_gap = (Inches(12.23)-4*top_cw)/3
    bot_gap = (Inches(12.23)-3*bot_cw)/2
    for i in range(4):
        cx = Inches(0.55)+i*(top_cw+top_gap)
        cy = top_y; cw = top_cw
        _card(sl,cx,cy,cw,card_h,pages[i])
    for i in range(3):
        cx = Inches(0.55)+i*(bot_cw+bot_gap)
        cy = bot_y; cw = bot_cw
        _card(sl,cx,cy,cw,card_h,pages[4+i])
    txt(sl,"http://localhost:3000  ·  docker-compose up",
        Inches(0.55),Inches(7.08),Inches(12.23),Inches(0.35),
        size=Pt(9.5),color=TEXT_M,italic=True,align=PP_ALIGN.CENTER)
    return sl

def _card(sl,cx,cy,cw,ch,page_data):
    icon,name,col,desc = page_data
    rrect(sl,cx,cy,cw,ch,fill=WHITE,line=BORDER)
    rect(sl,cx,cy,cw,Inches(0.065),fill=col)
    txt(sl,icon,cx+Inches(0.15),cy+Inches(0.1),Inches(0.5),Inches(0.45),size=Pt(18))
    txt(sl,name,cx+Inches(0.62),cy+Inches(0.16),cw-Inches(0.72),Inches(0.4),
        size=Pt(11),bold=True,color=TEXT_P)
    rect(sl,cx+Inches(0.15),cy+Inches(0.65),cw-Inches(0.3),Inches(0.014),fill=BORDER)
    txt(sl,desc,cx+Inches(0.15),cy+Inches(0.75),cw-Inches(0.3),Inches(1.75),
        size=Pt(9),color=TEXT_M,wrap=True)

# ── SLIDE 8 — Démonstration ──────────────────────────────────────────────────

def s8(prs):
    sl = new_slide(prs); bg(sl, BG)
    header(sl,"Démonstration — du budget au ROI en temps réel",
           "Cas réel : 90M€ TV · 10M€ Radio · 2M€ Social Media · Influenceur Mega.")
    steps = [
        ("①","Ajuster les sliders",TV,
         [("TV Budget","90M€",TV),("Radio Budget","10M€",RADIO),("Social Media","2M€",SOCIAL)],
         "Interface Framer Motion · sans clic requis"),
        ("②","Prédiction < 50 ms",SOCIAL,
         [("Revenue","$12.4M",GREEN),("ROI","+1140%",ACCENT),("Efficacité","11.2×",RADIO)],
         "Random Forest via API · ROI · efficacité · classe"),
        ("③","Analyser & Optimiser",GREEN,
         [("Rand. Forest","0.9965",TV),("Grad. Boost.","0.9964",RADIO),
          ("Lin. Reg.","0.9956",SOCIAL),("MLP","0.9941",INDIGO)],
         "4 modèles comparés · Planificateur budgétaire"),
    ]
    sw = Inches(3.9); sy = Inches(1.32)
    for i,(num,title,col,metrics,note) in enumerate(steps):
        sx = Inches(0.55)+i*Inches(4.26)
        sh = Inches(5.7)
        rrect(sl,sx,sy,sw,sh,fill=WHITE,line=BORDER)
        rect(sl,sx,sy,sw,Inches(0.075),fill=col)
        circ(sl,sx+Inches(0.2),sy+Inches(0.14),Inches(0.46),num,col)
        txt(sl,title,sx+Inches(0.82),sy+Inches(0.18),Inches(2.9),Inches(0.4),
            size=Pt(13),bold=True,color=TEXT_P)
        rect(sl,sx+Inches(0.15),sy+Inches(0.67),sw-Inches(0.3),Inches(0.014),fill=BORDER)
        mh = Inches(1.1) if len(metrics)==4 else Inches(1.35)
        for j,(ml,mv,mc) in enumerate(metrics):
            my = sy+Inches(0.82)+j*(mh+Inches(0.1))
            rrect(sl,sx+Inches(0.18),my,sw-Inches(0.36),mh,
                  fill=RGBColor(0xF8,0xFA,0xFF),line=BORDER,lw=Pt(0.3))
            rect(sl,sx+Inches(0.18),my,Inches(0.055),mh,fill=mc)
            txt(sl,ml,sx+Inches(0.35),my+Inches(0.06),Inches(2.5),Inches(0.28),
                size=Pt(8.5),color=TEXT_M)
            txt(sl,mv,sx+Inches(0.35),my+mh-Inches(0.46),Inches(2.8),Inches(0.38),
                size=Pt(15),bold=True,color=mc)
        ny = sy+sh-Inches(0.48)
        txt(sl,note,sx+Inches(0.18),ny,sw-Inches(0.36),Inches(0.4),
            size=Pt(8.5),color=TEXT_M,italic=True,wrap=True)
        if i<2:
            txt(sl,"›",sx+sw+Inches(0.14),sy+Inches(2.5),Inches(0.28),Inches(0.55),
                size=Pt(22),color=TEXT_M,align=PP_ALIGN.CENTER)
    return sl

# ── SLIDE 9 — Bilan critique ─────────────────────────────────────────────────

def s9(prs):
    sl = new_slide(prs); bg(sl, BG)
    header(sl,"Bilan critique",
           "Un système performant et démontrable aujourd'hui, avec une trajectoire claire d'amélioration.")
    cols_data = [
        ("💪  Forces",GREEN,[
            "R² > 0.99 sur les 4 modèles — prédictions ultra-précises",
            "Architecture modulaire et Docker-ready",
            "API REST documentée et scalable (FastAPI)",
            "Dashboard : 7 pages, mode clair/sombre, responsive",
            "Pipeline sklearn reproductible (splits.pkl, joblib)",
            "Déployable : docker-compose up",
        ]),
        ("⚠️  Limites",AMBER,[
            "Dataset synthétique — à valider sur données réelles",
            "Pas de feedback ni de ré-entraînement automatique",
            "Saisonnalité et tendances non modélisées",
            "Influencer : variable catégorielle simple",
            "Pas d'authentification sur l'API publique",
        ]),
        ("🚀  Perspectives",ACCENT,[
            "Intégration APIs réelles (Google Ads, Meta Ads)",
            "Modèles de séries temporelles (Prophet, LSTM)",
            "Ré-entraînement automatique (MLOps, Airflow)",
            "Attribution multi-touch (Shapley values)",
            "Déploiement cloud : GCP / AWS + auto-scaling",
        ]),
    ]
    cw,ch = Inches(3.98),Inches(5.7)
    cy = Inches(1.3)
    for i,(title,col,items) in enumerate(cols_data):
        cx = Inches(0.55)+i*Inches(4.25)
        rrect(sl,cx,cy,cw,ch,fill=WHITE,line=BORDER)
        rect(sl,cx,cy,cw,Inches(0.07),fill=col)
        rect(sl,cx,cy,Inches(0.07),ch,fill=col)
        txt(sl,title,cx+Inches(0.2),cy+Inches(0.12),cw-Inches(0.3),Inches(0.42),
            size=Pt(13),bold=True,color=col)
        rect(sl,cx+Inches(0.15),cy+Inches(0.62),cw-Inches(0.3),Inches(0.014),fill=BORDER)
        for j,item in enumerate(items):
            iy = cy+Inches(0.76)+j*Inches(0.8)
            oval(sl,cx+Inches(0.22),iy+Inches(0.15),Inches(0.09),Inches(0.09),fill=col)
            txt(sl,item,cx+Inches(0.42),iy,cw-Inches(0.58),Inches(0.7),
                size=Pt(9.5),color=TEXT_S,wrap=True)
    return sl

# ── SLIDE 10 — Conclusion ────────────────────────────────────────────────────

def s10(prs):
    sl = new_slide(prs); bg(sl, BG)
    rect(sl,0,0,W,Inches(0.09),fill=ACCENT)
    txt(sl,"Conclusion",Inches(0.55),Inches(0.18),Inches(12.0),Inches(0.55),
        size=Pt(25),bold=True,color=TEXT_P)
    rrect(sl,Inches(0.55),Inches(0.85),Inches(12.23),Inches(0.85),fill=BLUE_L,line=ACCENT)
    txt(sl,"ROI Intelligence démontre qu'il est possible de transformer un budget publicitaire\nen décision data-driven — précise, explicable et actionnable.",
        Inches(0.75),Inches(0.96),Inches(12.0),Inches(0.68),
        size=Pt(12.5),color=BLUE_D,italic=True)
    pillars = [
        ("⚡","Performance",   ACCENT,
         "Quatre modèles ML avec R² > 0.994, évalués par validation croisée et sélectionnés via RandomizedSearchCV."),
        ("🧱","Scalabilité",   RADIO,
         "Architecture Docker-ready : API FastAPI indépendante du frontend, extensible à de nouveaux canaux."),
        ("🔍","Explicabilité", GREEN,
         "Dashboard 7 pages : importance des features, sensibilité par canal, planification inverse — tout est justifié."),
    ]
    pw,ph = Inches(3.87),Inches(2.55)
    py = Inches(1.9)
    for i,(icon,title,col,desc) in enumerate(pillars):
        px = Inches(0.55)+i*Inches(4.27)
        rrect(sl,px,py,pw,ph,fill=WHITE,line=BORDER)
        rect(sl,px,py,pw,Inches(0.07),fill=col)
        ic = oval(sl,px+Inches(0.2),py+Inches(0.18),Inches(0.5),Inches(0.5),fill=col)
        stxt(ic,icon,size=Pt(16),bold=False)
        txt(sl,title,px+Inches(0.86),py+Inches(0.21),Inches(2.8),Inches(0.4),
            size=Pt(14),bold=True,color=TEXT_P)
        rect(sl,px+Inches(0.15),py+Inches(0.76),pw-Inches(0.3),Inches(0.014),fill=BORDER)
        txt(sl,desc,px+Inches(0.15),py+Inches(0.9),pw-Inches(0.3),Inches(1.48),
            size=Pt(10.5),color=TEXT_S,wrap=True)
    rect(sl,Inches(0.55),Inches(4.65),Inches(12.23),Inches(0.018),fill=BORDER)
    txt(sl,"Merci pour votre attention.",
        Inches(0.55),Inches(4.8),Inches(12.23),Inches(0.55),
        size=Pt(22),bold=True,color=TEXT_P,align=PP_ALIGN.CENTER)
    txt(sl,"Nous sommes disponibles pour répondre à vos questions.",
        Inches(0.55),Inches(5.4),Inches(12.23),Inches(0.4),
        size=Pt(13),color=TEXT_M,align=PP_ALIGN.CENTER)
    txt(sl,"CHUEMBOU ADRIEN · QUANG DAT  —  EFREI Paris · M1 Data Science & IA · 2025–2026",
        Inches(0.55),Inches(6.9),Inches(12.23),Inches(0.38),
        size=Pt(10),color=TEXT_M,align=PP_ALIGN.CENTER)
    return sl

# ── BUILD ────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    s1(prs); s2(prs); s3(prs); s4(prs); s5(prs)
    s6(prs); s7(prs); s8(prs); s9(prs); s10(prs)
    out = "ROI_Intelligence_Presentation_v2.pptx"
    prs.save(out)
    print(f"Saved: {out}")

build()
