#!/usr/bin/env python3
"""Generate ROI Intelligence presentation matching the dashboard design system."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SH

# ── Design tokens ──────────────────────────────────────────────────────────────
BG      = RGBColor(0x0e, 0x15, 0x25)
CARD    = RGBColor(0x0f, 0x17, 0x2a)
CARD2   = RGBColor(0x12, 0x1c, 0x30)
BORDER  = RGBColor(0x1e, 0x29, 0x3b)
ACCENT  = RGBColor(0x3b, 0x82, 0xf6)
PURPLE  = RGBColor(0x8b, 0x5c, 0xf6)
CYAN    = RGBColor(0x06, 0xb6, 0xd4)
GREEN   = RGBColor(0x10, 0xb9, 0x81)
ORANGE  = RGBColor(0xf5, 0x9e, 0x0b)
RED     = RGBColor(0xef, 0x44, 0x44)
PINK    = RGBColor(0xec, 0x48, 0x99)
TEAL    = RGBColor(0x14, 0xb8, 0xa6)
WHITE   = RGBColor(0xe2, 0xe8, 0xf0)
MUTED   = RGBColor(0x64, 0x74, 0x8b)
SEC     = RGBColor(0x94, 0xa3, 0xb8)
FONT    = "Calibri"

SW = 13.333
SH_ = 7.5

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH_)
    return prs

def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def rect(slide, x, y, w, h, fill=CARD, line=None, lw=0.5):
    sh = slide.shapes.add_shape(SH.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else:    sh.line.fill.background()
    return sh

def card(slide, x, y, w, h, fill=CARD, line=BORDER, lw=0.5):
    sh = slide.shapes.add_shape(SH.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else:    sh.line.fill.background()
    return sh

def oval(slide, x, y, d, fill=ACCENT):
    sh = slide.shapes.add_shape(SH.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    return sh

def txt(slide, text, x, y, w, h, sz=12, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color
    return tb

def icon_circle(slide, cx, cy, d, icon_text, icon_color):
    o = oval(slide, cx, cy, d, fill=icon_color)
    tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = icon_text
    r.font.size = Pt(int(d * 30)); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = FONT

def header(slide, title, sub=None, eyebrow=None):
    y = 0.32
    if eyebrow:
        txt(slide, eyebrow.upper(), 0.5, y, 12, 0.3, sz=8, bold=True, color=ACCENT)
        y += 0.27
    txt(slide, title, 0.5, y, 12.3, 0.62, sz=26, bold=True, color=WHITE)
    y += 0.58
    if sub:
        txt(slide, sub, 0.5, y, 12.3, 0.42, sz=12, italic=True, color=MUTED)
        y += 0.38
    rect(slide, 0.5, y, 1.4, 0.04, fill=ACCENT)
    return y + 0.18

# ─────────────────────────────────────────────────────────────────────────────
def s01_title(prs):
    s = blank(prs)
    # dot grid background
    for c in range(22):
        for r_ in range(12):
            d = oval(s, c*0.63-0.2, r_*0.67-0.1, 0.04, fill=RGBColor(0x1a,0x26,0x3a))
    # left accent bar
    rect(s, 0.45, 1.75, 0.09, 2.6, fill=ACCENT)
    # title
    txt(s, "ROI Intelligence", 0.75, 1.72, 9, 1.1, sz=52, bold=True, color=WHITE)
    rect(s, 0.75, 2.8, 5.8, 0.055, fill=ACCENT)
    txt(s, "Optimisation du retour sur investissement marketing\npar le Machine Learning",
        0.75, 2.95, 9, 0.85, sz=19, color=ACCENT)
    txt(s, "Système de prédiction multi-modèles et tableau de bord décisionnel en temps réel",
        0.75, 3.82, 9, 0.45, sz=12, italic=True, color=MUTED)
    rect(s, 0.5, 5.1, 12.3, 0.015, fill=BORDER)
    txt(s, "Soutenance · EFREI Paris · M1 Data Science & IA · 2025–2026",
        0.5, 5.25, 9, 0.35, sz=10, color=MUTED)
    txt(s, "CHUEMBOU ADRIEN  ·  QUANG DAT", 0.5, 5.62, 9, 0.38, sz=12, bold=True, color=SEC)
    # right metric cards
    items = [("R²", "Score modèle", ACCENT), ("×4", "Modèles ML", PURPLE),
             ("×7", "Pages dash.", CYAN),    ("<50ms","API latence", GREEN)]
    for i,(v,lbl,c) in enumerate(items):
        cx = 10.2 + (i%2)*1.55; cy = 1.8 + (i//2)*1.32
        card(s, cx, cy, 1.38, 1.18, fill=CARD2, line=c, lw=0.9)
        txt(s, v,   cx+0.06, cy+0.12, 1.26, 0.55, sz=22, bold=True, color=c, align=PP_ALIGN.CENTER)
        txt(s, lbl, cx+0.06, cy+0.68, 1.26, 0.35, sz=9,  color=MUTED, align=PP_ALIGN.CENTER)
    return s

# ─────────────────────────────────────────────────────────────────────────────
def s02_problematique(prs):
    s = blank(prs)
    cy = header(s, "Problématique",
                "Les budgets publicitaires sont encore trop souvent alloués à l'intuition, pas à la donnée.")
    cw, ch = 5.85, SH_ - cy - 0.45
    card(s, 0.45, cy, cw, ch, fill=RGBColor(0x14,0x0d,0x0d), line=RED, lw=0.6)
    card(s, 0.45+cw+0.28, cy, cw, ch, fill=RGBColor(0x0a,0x13,0x25), line=ACCENT, lw=0.6)
    rx = 0.45+cw+0.28
    txt(s, "❌  Approche traditionnelle", 0.65, cy+0.15, cw-0.35, 0.4, sz=13, bold=True, color=RED)
    txt(s, "✅  Notre approche ML",       rx+0.2,cy+0.15, cw-0.35, 0.4, sz=13, bold=True, color=GREEN)
    rect(s, 0.65, cy+0.55, cw-0.35, 0.015, fill=RGBColor(0x3a,0x1a,0x1a))
    rect(s, rx+0.2,cy+0.55, cw-0.35, 0.015, fill=BORDER)
    L = ["Allocation empirique basée sur l'expérience passée",
         "Aucune corrélation mesurée entre budget et revenu",
         "Décisions prises a posteriori, après les campagnes",
         "Même budget, résultats imprévisibles selon les canaux",
         "Impossible de simuler l'impact avant d'engager"]
    R = ["Modèles entraînés sur 4 572 campagnes documentées",
         "Prédiction du CA en temps réel selon le budget alloué",
         "4 canaux : TV · Radio · Social Media · Influenceur",
         "R² > 0.994 : corrélation quasi-parfaite avec les ventes",
         "Simulation live : le revenu s'affiche instantanément"]
    for i,(l,r_) in enumerate(zip(L,R)):
        by = cy+0.72+i*0.54
        oval(s, 0.67, by+0.1, 0.07, fill=RED)
        txt(s, l, 0.86, by, cw-0.55, 0.48, sz=11, color=SEC)
        oval(s, rx+0.22, by+0.1, 0.07, fill=GREEN)
        txt(s, r_, rx+0.41, by, cw-0.55, 0.48, sz=11, color=SEC)
    rect(s, 0.45+cw+0.1, cy+0.2, 0.08, ch-0.4, fill=BORDER)
    qy = cy+ch+0.08
    txt(s, "L'enjeu : transformer une intuition budgétaire en décision data-driven, mesurable et reproductible.",
        0.5, qy, 12.3, 0.38, sz=11, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    return s

# ─────────────────────────────────────────────────────────────────────────────
def s03_objectives(prs):
    s = blank(prs)
    cy = header(s, "Objectifs du projet",
                "Quatre ambitions concrètes, de la donnée brute à la décision business éclairée.")
    cards = [
        ("①","Prédire",ACCENT,
         "Construire 4 modèles ML capables de prédire le chiffre d'affaires\nà partir du mix budgétaire TV · Radio · Social Media · Influenceur."),
        ("②","Comparer & Sélectionner",PURPLE,
         "Évaluer chaque modèle via validation croisée 5 plis (R², MAE, RMSE)\net sélectionner automatiquement le plus performant."),
        ("③","Exposer via API REST",CYAN,
         "Déployer les modèles entraînés via une API FastAPI scalable,\ninterrogeable en temps réel par n'importe quel client."),
        ("④","Visualiser & Décider",GREEN,
         "Proposer un dashboard interactif en 7 pages permettant\nde simuler, comparer et optimiser les allocations budgétaires."),
    ]
    cw, ch, gx, gy = 5.9, 2.42, 0.28, 0.22
    for i,(num,title,color,desc) in enumerate(cards):
        col,row = i%2, i//2
        cx = 0.48+col*(cw+gx); cy2 = cy+row*(ch+gy)
        card(s, cx, cy2, cw, ch, fill=CARD, line=color, lw=0.9)
        icon_circle(s, cx+0.2, cy2+0.18, 0.58, num, color)
        txt(s, title, cx+0.9, cy2+0.2, cw-1.05, 0.45, sz=15, bold=True, color=WHITE)
        rect(s, cx+0.2, cy2+0.72, cw-0.4, 0.02, fill=BORDER)
        txt(s, desc, cx+0.2, cy2+0.82, cw-0.4, 1.45, sz=11, color=SEC)
    return s

# ─────────────────────────────────────────────────────────────────────────────
def s04_architecture(prs):
    s = blank(prs)
    cy = header(s, "Architecture de la solution",
                "Cinq couches modulaires : de la donnée brute à la décision business.")
    layers = [
        ("DONNÉES",        "data/Dummy Data HSS.csv",
         "4 572 lignes · TV, Radio, Social Media, Influencer, Sales", CYAN),
        ("PREPROCESSING",  "src/preprocessing.py",
         "StandardScaler + OneHotEncoder · Pipeline sklearn · Imputation médiane", PURPLE),
        ("ENTRAÎNEMENT ML","src/train.py",
         "RandomizedSearchCV · 15 iter · 5-fold CV · 4 algorithmes → .pkl (joblib)", ACCENT),
        ("API REST",       "api/main.py : FastAPI",
         "/predict · /predict/all · /metrics · /feature-importance · /stats · /health", GREEN),
        ("DASHBOARD",      "dashboard-next/ : Next.js 14",
         "TypeScript · Tailwind CSS · Recharts · Framer Motion · 7 pages · Clair/Sombre", ORANGE),
    ]
    lh, lw, gap = 0.7, 12.0, 0.13
    for i,(lbl,title,desc,color) in enumerate(layers):
        ly = cy + i*(lh+gap)
        card(s, 0.5, ly, lw, lh, fill=CARD, line=color, lw=0.7)
        rect(s, 0.5, ly, 0.11, lh, fill=color)
        # badge
        badge = s.shapes.add_shape(SH.RECTANGLE, Inches(0.74), Inches(ly+0.16), Inches(1.55), Inches(0.34))
        badge.fill.solid(); badge.fill.fore_color.rgb = color; badge.line.fill.background()
        btf = badge.text_frame; btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = lbl; br.font.size = Pt(8); br.font.bold = True
        br.font.color.rgb = WHITE; br.font.name = FONT
        txt(s, title, 2.5, ly+0.1, 5.5, 0.35, sz=12, bold=True, color=WHITE)
        txt(s, desc,  2.5, ly+0.4, 9.6, 0.28, sz=10, color=MUTED)
        if i < len(layers)-1:
            rect(s, 6.28, ly+lh+0.01, 0.07, gap-0.01, fill=ACCENT)
    note_y = cy + len(layers)*(lh+gap) + 0.06
    txt(s, "⚡  docker-compose up  ·  API :8000  ·  Dashboard :3000",
        0.5, note_y, 12.3, 0.33, sz=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    return s

# ─────────────────────────────────────────────────────────────────────────────
def s05_pipeline(prs):
    s = blank(prs)
    cy = header(s, "Pipeline ML : de la donnée brute à la prédiction",
                "Un pipeline rigoureusement validé : de la saisie brute à l'inférence en moins de 50 ms.")
    steps = [
        ("①","load_data()","CSV · 4 572 lignes\nValidation schéma",CYAN),
        ("②","build_preprocessor()","StandardScaler\n+ OneHotEncoder",PURPLE),
        ("③","tune()","RandomizedSearchCV\n15 iter · 5-fold",ACCENT),
        ("④","train_all()","4 modèles .pkl\nsplits.pkl",GREEN),
        ("⑤","FastAPI","Lifespan loading\n< 50ms / req.",ORANGE),
        ("⑥","Dashboard","Sliders → API\nTemps réel",PINK),
    ]
    sw2, sh2, gx = 1.92, 2.25, 0.1
    total = len(steps)*sw2+(len(steps)-1)*gx
    sx0 = (SW-total)/2
    sy = cy+0.15
    for i,(num,title,desc,color) in enumerate(steps):
        sx = sx0+i*(sw2+gx)
        card(s, sx, sy, sw2, sh2, fill=CARD, line=color, lw=0.9)
        icon_circle(s, sx+sw2/2-0.29, sy+0.16, 0.58, num, color)
        txt(s, title, sx+0.08, sy+0.84, sw2-0.16, 0.52, sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(s, sx+0.15, sy+1.38, sw2-0.3, 0.02, fill=BORDER)
        txt(s, desc,  sx+0.1,  sy+1.48, sw2-0.2, 0.68, sz=10, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(steps)-1:
            txt(s, "›", sx+sw2+0.01, sy+sh2/2-0.22, gx+0.06, 0.4, sz=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    note_y = sy+sh2+0.2
    card(s, 1.4, note_y, 10.5, 0.48, fill=CARD, line=ACCENT, lw=0.5)
    txt(s, "🔒  models/splits.pkl garantit la reproductibilité : métriques toujours calculées sur le même jeu de test.",
        1.6, note_y+0.09, 10.1, 0.32, sz=10, italic=True, color=MUTED)
    return s

# ─────────────────────────────────────────────────────────────────────────────
def s06_models(prs):
    s = blank(prs)
    cy = header(s, "Les 4 modèles ML : Performances comparées",
                "Tous les modèles dépassent R² = 0.99 : le Random Forest s'impose comme référence.")
    # Table
    tx, tw = 0.5, 7.1
    cols  = ["Modèle","R² moyen","Écart-type","Statut"]
    cws   = [2.9, 1.6, 1.5, 1.1]
    rows  = [
        ("🏆  Random Forest",    "0.9965","±0.0028","BEST", ACCENT, True),
        ("Gradient Boosting",    "0.9964","±0.0028","—",    SEC,    False),
        ("Régression Linéaire",  "0.9956","±0.0032","—",    SEC,    False),
        ("MLP Neural Network",   "0.9941","±0.0034","—",    SEC,    False),
    ]
    hh, rh = 0.42, 0.52
    rect(s, tx, cy, tw, hh, fill=RGBColor(0x1e,0x2d,0x48))
    hx = tx
    for hdr,cw in zip(cols,cws):
        txt(s, hdr, hx+0.1, cy+0.06, cw-0.1, hh-0.1, sz=11, bold=True, color=ACCENT)
        hx += cw
    for i,(model,r2,std,status,color,best) in enumerate(rows):
        ry = cy+hh+i*rh
        rf = RGBColor(0x12,0x20,0x3a) if best else (CARD if i%2==0 else RGBColor(0x0d,0x14,0x24))
        rect(s, tx, ry, tw, rh, fill=rf)
        if best: rect(s, tx, ry, 0.08, rh, fill=ACCENT)
        vals = [model,r2,std,status]
        cx2 = tx
        for j,(v,cw) in enumerate(zip(vals,cws)):
            vc = (ACCENT if best else SEC) if j>0 else (WHITE if best else SEC)
            txt(s, v, cx2+0.15, ry+0.1, cw-0.15, rh-0.1, sz=12, bold=best, color=vc)
            cx2 += cw
        rect(s, tx, ry+rh-0.01, tw, 0.01, fill=BORDER)
    # Bar chart
    bx, bw = 8.0, 5.0
    txt(s, "R² par modèle (5-fold CV)", bx, cy, bw, 0.38, sz=11, bold=True, color=SEC)
    bms = [("Random Forest",0.9965,ACCENT),("Gradient Boosting",0.9964,PURPLE),
           ("Régression Lin.",0.9956,CYAN),  ("MLP Neural Net",  0.9941,MUTED)]
    mn, mx2 = 0.993, 0.997
    for i,(name,val,color) in enumerate(bms):
        by = cy+0.55+i*0.6
        bw2 = bw-1.3
        rect(s, bx, by+0.08, bw2, 0.32, fill=BORDER)
        frac = (val-mn)/(mx2-mn)
        rect(s, bx, by+0.08, max(0.08,bw2*frac), 0.32, fill=color if i==0 else RGBColor(0x30,0x45,0x65))
        txt(s, name, bx, by, bw2, 0.26, sz=9, color=SEC)
        txt(s, str(val), bx+bw2+0.08, by+0.06, 0.88, 0.32, sz=11, bold=(i==0), color=color)
    # Highlight card
    hy = cy+3.15
    card(s, bx, hy, bw, 1.45, fill=RGBColor(0x0a,0x16,0x30), line=ACCENT, lw=1.0)
    txt(s, "🏆  Meilleur modèle", bx+0.2, hy+0.1, bw-0.4, 0.32, sz=11, bold=True, color=ACCENT)
    txt(s, "Random Forest",       bx+0.2, hy+0.42, bw-0.4, 0.48, sz=22, bold=True, color=WHITE)
    txt(s, "R² = 0.9965  ·  Explique 99.65% de la variance des ventes",
        bx+0.2, hy+0.92, bw-0.4, 0.38, sz=10, italic=True, color=MUTED)
    return s

# ─────────────────────────────────────────────────────────────────────────────
def s07_dashboard(prs):
    s = blank(prs)
    cy = header(s, "Dashboard : 7 pages d'intelligence marketing",
                "De la prédiction brute à l'optimisation stratégique : chaque page répond à une question business précise.")
    pages = [
        ("📊","Accueil",ACCENT,         "Vue globale · KPIs · corrélation canaux · benchmark revenus"),
        ("🎯","Prévision des revenus",PURPLE,"Sliders · prédiction temps réel · ROI · consensus 4 modèles"),
        ("⚙️","Optimiseur budgétaire",CYAN, "Présets · sensibilité TV · comparaison des modèles"),
        ("🗺️","Planificateur cible",GREEN,  "Objectif → budget optimal · inversion modèle · Monte Carlo RF"),
        ("🤖","Comparaison modèles",ORANGE, "Métriques · diagramme radar · matrice de confusion"),
        ("💡","Insights & Corrélations",PINK,"Importance features · corrélation · ROI par influenceur"),
        ("📈","Feature Importance",TEAL,   "Classement SHAP · permutation importance · stats canal-revenu"),
    ]
    cw, ch, gx, gy = 3.88, 1.35, 0.19, 0.2
    for i,(icon,name,color,desc) in enumerate(pages):
        col,row = i%3, i//3
        cx = (SW-cw)/2 if i==6 else 0.48+col*(cw+gx)
        cy2 = cy+row*(ch+gy)
        card(s, cx, cy2, cw, ch, fill=CARD, line=color, lw=0.6)
        rect(s, cx, cy2, cw, 0.07, fill=color)
        txt(s, icon+"  "+name, cx+0.15, cy2+0.14, cw-0.3, 0.4, sz=12, bold=True, color=WHITE)
        rect(s, cx+0.15, cy2+0.56, cw-0.3, 0.015, fill=BORDER)
        txt(s, desc, cx+0.15, cy2+0.64, cw-0.3, 0.6, sz=10, color=MUTED)
    return s

# ─────────────────────────────────────────────────────────────────────────────
def s08_demo(prs):
    s = blank(prs)
    cy = header(s, "Démonstration : du budget au ROI en temps réel",
                "Cas réel : 90M€ TV · 10M€ Radio · 2M€ Social Media · Influenceur Mega.")
    steps = [
        ("①","Ajuster les sliders",ACCENT,"Interface Framer Motion\naucun bouton requis",
         "sliders"),
        ("②","Prédiction < 50 ms",GREEN, "Random Forest via API\nROI · efficacité · classe",
         "metrics"),
        ("③","Analyser & Optimiser",ORANGE,"4 modèles comparés\nPlanificateur budgétaire",
         "bars"),
    ]
    sw2, sh2, gx = 4.05, 4.6, 0.25
    sx0 = (SW - (len(steps)*sw2+(len(steps)-1)*gx)) / 2
    for i,(num,title,color,desc,viz) in enumerate(steps):
        sx = sx0+i*(sw2+gx)
        card(s, sx, cy, sw2, sh2, fill=CARD, line=color, lw=0.9)
        icon_circle(s, sx+0.2, cy+0.18, 0.56, num, color)
        txt(s, title, sx+0.9, cy+0.22, sw2-1.05, 0.4, sz=14, bold=True, color=WHITE)
        txt(s, desc,  sx+0.9, cy+0.62, sw2-1.05, 0.42, sz=10, color=color)
        # mockup area
        my = cy+1.12; mh = 2.0
        card(s, sx+0.18, my, sw2-0.36, mh, fill=RGBColor(0x07,0x0f,0x1c), line=BORDER, lw=0.3)
        if viz == "sliders":
            lbls = ["TV Budget","Radio Budget","Social Media"]
            vals = [0.90, 0.10, 0.04]
            vcols = [ACCENT, PURPLE, CYAN]
            vlbls = ["90M€","10M€","2M€"]
            for j,(lb,v,c,vl) in enumerate(zip(lbls,vals,vcols,vlbls)):
                by = my+0.2+j*0.59
                txt(s, lb, sx+0.32, by, 1.5, 0.24, sz=9, color=MUTED)
                bx0 = sx+0.32; bw3 = sw2-0.68
                rect(s, bx0, by+0.27, bw3, 0.14, fill=BORDER)
                rect(s, bx0, by+0.27, bw3*v, 0.14, fill=c)
                oval(s, bx0+bw3*v-0.1, by+0.22, 0.2, fill=WHITE)
                txt(s, vl, bx0+bw3-0.65, by, 0.6, 0.24, sz=9, bold=True, color=c, align=PP_ALIGN.RIGHT)
        elif viz == "metrics":
            mdata = [("Revenue","$12.4M",ACCENT),("ROI","+1140%",GREEN),("Eff.","11.2×",PURPLE)]
            for j,(lb,v,c) in enumerate(mdata):
                mx2 = sx+0.28+j*1.18
                card(s, mx2, my+0.2, 1.06, 1.45, fill=CARD, line=c, lw=0.5)
                txt(s, lb, mx2+0.05, my+0.3,  1.0, 0.28, sz=8,  color=MUTED, align=PP_ALIGN.CENTER)
                txt(s, v,  mx2+0.05, my+0.62, 1.0, 0.52, sz=14, bold=True, color=c, align=PP_ALIGN.CENTER)
        else:
            bdata = [("Rand. Forest",0.98,ACCENT),("Grad. Boost.",0.96,PURPLE),
                     ("Lin. Reg.",0.88,CYAN),     ("MLP",0.82,MUTED)]
            for j,(lb,v,c) in enumerate(bdata):
                by2 = my+0.22+j*0.44
                bw4 = sw2-0.72
                rect(s, sx+0.3, by2+0.06, bw4, 0.26, fill=BORDER)
                rect(s, sx+0.3, by2+0.06, bw4*v, 0.26, fill=c)
                txt(s, lb, sx+0.3, by2, bw4*0.6, 0.22, sz=8, color=MUTED)
                txt(s, f"{v:.2f}", sx+0.3+bw4+0.04, by2+0.04, 0.45, 0.26, sz=9, bold=True, color=c)
        # description
        rect(s, sx+0.18, my+mh+0.06, sw2-0.36, 0.015, fill=BORDER)
        txt(s, desc, sx+0.18, my+mh+0.12, sw2-0.36, 1.1, sz=10, color=MUTED)
    return s

# ─────────────────────────────────────────────────────────────────────────────
def s09_bilan(prs):
    s = blank(prs)
    cy = header(s, "Bilan critique",
                "Un système performant et démontrable aujourd'hui, avec une trajectoire claire d'amélioration.")
    cols_data = [
        ("💪","Forces",GREEN,
         ["R² > 0.99 sur les 4 modèles : prédictions ultra-précises",
          "Architecture modulaire et Docker-ready",
          "API REST documentée et scalable (FastAPI)",
          "Dashboard : 7 pages, mode clair/sombre, responsive",
          "Pipeline sklearn reproductible (splits.pkl, joblib)",
          "Déployable : docker-compose up"]),
        ("⚠️","Limites",ORANGE,
         ["Dataset synthétique : à valider sur données réelles",
          "Pas de feedback ni de ré-entraînement automatique",
          "Saisonnalité et tendances non modélisées",
          "Influencer : variable catégorielle simple",
          "Pas d'authentification sur l'API publique"]),
        ("🚀","Perspectives",ACCENT,
         ["Intégration APIs réelles (Google Ads, Meta Ads)",
          "Modèles de séries temporelles (Prophet, LSTM)",
          "Ré-entraînement automatique (MLOps, Airflow)",
          "Attribution multi-touch (Shapley values)",
          "Déploiement cloud : GCP / AWS + auto-scaling"]),
    ]
    cw, ch, gx = 3.9, SH_-cy-0.35, 0.27
    for i,(icon,title,color,bullets) in enumerate(cols_data):
        cx = 0.48+i*(cw+gx)
        card(s, cx, cy, cw, ch, fill=CARD, line=color, lw=0.55)
        rect(s, cx, cy, cw, 0.09, fill=color)
        txt(s, icon+"  "+title, cx+0.2, cy+0.18, cw-0.4, 0.5, sz=17, bold=True, color=color)
        rect(s, cx+0.2, cy+0.7, cw-0.4, 0.02, fill=color)
        for j,bullet in enumerate(bullets):
            by = cy+0.85+j*0.67
            oval(s, cx+0.22, by+0.1, 0.08, fill=color)
            txt(s, bullet, cx+0.42, by, cw-0.55, 0.55, sz=11, color=SEC)
    return s

# ─────────────────────────────────────────────────────────────────────────────
def s10_conclusion(prs):
    s = blank(prs)
    oval(s, 3.8, 0.4, 5.8, fill=RGBColor(0x0a,0x14,0x2e))
    txt(s, "Conclusion", 0.5, 0.35, 12.3, 0.72, sz=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 5.6, 1.1, 2.1, 0.055, fill=ACCENT)
    txt(s, "ROI Intelligence démontre qu'il est possible de transformer un budget publicitaire\nen décision data-driven : précise, explicable et actionnable.",
        0.8, 1.28, 11.7, 0.78, sz=13, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    pillars = [
        ("⚡","Performance",ACCENT,
         "Quatre modèles ML avec R² > 0.994, évalués par validation\ncroisée et sélectionnés via RandomizedSearchCV."),
        ("🧱","Scalabilité",CYAN,
         "Architecture Docker-ready : API FastAPI indépendante\ndu frontend, extensible à de nouveaux canaux."),
        ("🔍","Explicabilité",PURPLE,
         "Dashboard 7 pages : importance des features, sensibilité\npar canal, planification inverse, tout est justifié."),
    ]
    pw, ph, gx, py = 3.8, 2.42, 0.27, 2.3
    for i,(icon,title,color,desc) in enumerate(pillars):
        px = 0.6+i*(pw+gx)
        card(s, px, py, pw, ph, fill=RGBColor(0x0d,0x18,0x2e), line=color, lw=1.0)
        icon_circle(s, px+pw/2-0.33, py+0.18, 0.66, icon, color)
        txt(s, title, px+0.15, py+1.0, pw-0.3, 0.45, sz=16, bold=True, color=color, align=PP_ALIGN.CENTER)
        rect(s, px+0.4, py+1.5, pw-0.8, 0.02, fill=color)
        txt(s, desc, px+0.15, py+1.6, pw-0.3, 0.7, sz=11, color=MUTED, align=PP_ALIGN.CENTER)
    sep_y = py+ph+0.22
    rect(s, 0.5, sep_y, 12.3, 0.015, fill=BORDER)
    txt(s, "Merci pour votre attention.", 0.5, sep_y+0.1, 12.3, 0.42,
        sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "Nous sommes disponibles pour répondre à vos questions.",
        0.5, sep_y+0.54, 12.3, 0.35, sz=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    txt(s, "CHUEMBOU ADRIEN · QUANG DAT · EFREI Paris · M1 Data Science & IA · 2025–2026",
        0.5, sep_y+0.9, 12.3, 0.32, sz=10, color=MUTED, align=PP_ALIGN.CENTER)
    return s

# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = new_prs()
    s01_title(prs)
    s02_problematique(prs)
    s03_objectives(prs)
    s04_architecture(prs)
    s05_pipeline(prs)
    s06_models(prs)
    s07_dashboard(prs)
    s08_demo(prs)
    s09_bilan(prs)
    s10_conclusion(prs)
    out = "ROI_Intelligence_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")
